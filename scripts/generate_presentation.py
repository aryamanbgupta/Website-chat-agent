"""
Generate PartSelect AI Chat Agent presentation (PPTX).

Usage:
    python scripts/generate_presentation.py

Output:
    PartSelect_AI_Agent_Presentation.pptx (in project root)
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ──────────────────────────────────────────────
# Constants
# ──────────────────────────────────────────────
PROJECT_ROOT = Path(__file__).resolve().parent.parent
SCREENSHOTS_DIR = PROJECT_ROOT / "docs" / "screenshots"
OUTPUT_PATH = PROJECT_ROOT / "PartSelect_AI_Agent_Presentation.pptx"

# Colors
TEAL = RGBColor(0x33, 0x77, 0x78)
YELLOW = RGBColor(0xF3, 0xC0, 0x4C)
RED = RGBColor(0xAA, 0x1E, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
TABLE_HEADER_BG = TEAL
TABLE_ALT_ROW = RGBColor(0xF0, 0xF7, 0xF7)

FONT_NAME = "Calibri"
TITLE_SIZE = Pt(32)
SUBTITLE_SIZE = Pt(20)
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(14)
CAPTION_SIZE = Pt(13)
CODE_SIZE = Pt(12)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FOOTER_TEXT = "PartSelect AI Chat Agent — Aryaman Gupta | Instalily Case Study"

# Title bar geometry
BAR_HEIGHT = Inches(1.05)
BAR_TOP = Inches(0)

# Content area
CONTENT_LEFT = Inches(0.7)
CONTENT_TOP = Inches(1.35)
CONTENT_WIDTH = Inches(11.9)


# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────
def _set_font(run, size=BODY_SIZE, bold=False, color=DARK, name=FONT_NAME, italic=False):
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    run.font.italic = italic


def add_title_bar(slide, title_text):
    """Add a teal bar at the top with white title text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), BAR_TOP, SLIDE_W, BAR_HEIGHT,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = title_text
    _set_font(run, size=TITLE_SIZE, bold=True, color=WHITE)
    # vertical centering
    tf.margin_left = Inches(0.7)
    tf.margin_top = Inches(0.15)


def add_footer(slide):
    """Add a small footer at the bottom-right."""
    left = Inches(6.5)
    top = Inches(7.05)
    width = Inches(6.5)
    height = Inches(0.35)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = FOOTER_TEXT
    _set_font(run, size=Pt(10), color=MID_GRAY)


def add_textbox(slide, text, left, top, width, height,
                size=BODY_SIZE, bold=False, color=DARK, alignment=PP_ALIGN.LEFT,
                italic=False, name=FONT_NAME):
    """Add a single-paragraph textbox and return the textframe."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color, name=name, italic=italic)
    return tf


def add_bullet_list(slide, items, left, top, width, height,
                    size=BODY_SIZE, color=DARK, bold_prefix=False, spacing=Pt(8)):
    """Add a bulleted list. Items can be str or (bold_part, rest) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        p.level = 0

        if isinstance(item, tuple):
            # (bold_header, rest_text)
            run_b = p.add_run()
            run_b.text = item[0]
            _set_font(run_b, size=size, bold=True, color=color)
            run_r = p.add_run()
            run_r.text = item[1]
            _set_font(run_r, size=size, color=color)
        else:
            bullet_char = "\u2022  "
            run = p.add_run()
            run.text = bullet_char + item
            _set_font(run, size=size, color=color)
    return tf


def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.45)):
    """Add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows)
    table = table_shape.table

    # Distribute column widths
    col_width = int(width / n_cols)
    for col_idx in range(n_cols):
        table.columns[col_idx].width = col_width

    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                _set_font(run, size=SMALL_SIZE, bold=True, color=WHITE)

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_ROW
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    _set_font(run, size=SMALL_SIZE, color=DARK)

    return table_shape


def add_image_or_placeholder(slide, filename, left, top, width, height):
    """Add image from screenshots dir, or a labeled gray placeholder."""
    img_path = SCREENSHOTS_DIR / filename
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), left, top, width, height)
    else:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = LIGHT_GRAY
        rect.line.color.rgb = MID_GRAY
        rect.line.width = Pt(1)
        tf = rect.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = f"[{filename}]"
        _set_font(run, size=SMALL_SIZE, color=MID_GRAY, italic=True)
        # vertical center
        tf.margin_top = int(height / 2) - Inches(0.2)


def add_code_block(slide, text, left, top, width, height):
    """Add a monospaced code block with light background."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    rect.line.color.rgb = LIGHT_GRAY
    rect.line.width = Pt(1)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.2)

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        _set_font(run, size=CODE_SIZE, color=DARK, name="Courier New")
    return tf


def add_callout_box(slide, text, left, top, width, height, bg_color=None):
    """Add a colored callout box with text."""
    if bg_color is None:
        bg_color = RGBColor(0xFD, 0xF5, 0xE6)  # light yellow
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color
    rect.line.fill.background()
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    _set_font(run, size=SMALL_SIZE, color=DARK, italic=True)
    return tf


def new_slide(prs):
    """Add a blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])


# ──────────────────────────────────────────────
# Slide builders
# ──────────────────────────────────────────────

def slide_01_title(prs):
    slide = new_slide(prs)
    # Teal accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    # Title
    add_textbox(slide, "PartSelect AI Chat Agent",
                Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
                size=Pt(44), bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, "An AI-powered parts assistant built with a custom agent loop",
                Inches(1.5), Inches(3.3), Inches(10), Inches(0.8),
                size=SUBTITLE_SIZE, color=DARK, alignment=PP_ALIGN.CENTER)

    # Yellow accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = YELLOW
    line.line.fill.background()

    # Bottom
    add_textbox(slide, "Aryaman Gupta  \u00b7  Case Study for Instalily AI  \u00b7  March 2026",
                Inches(1.5), Inches(5.2), Inches(10), Inches(0.6),
                size=SMALL_SIZE, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_02_problem(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "The Problem")
    add_footer(slide)

    items = [
        "PartSelect.com has 2M+ OEM appliance parts — overwhelming for customers",
        "No AI-assisted search or diagnosis today — only basic live chat and phone support",
        "Customers struggle to: find the right replacement part, verify compatibility with their model, diagnose symptoms, and understand installation difficulty",
    ]
    add_bullet_list(slide, items, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(4.0),
                    size=BODY_SIZE)

    add_textbox(slide,
                "42 candidates forked the case study repo. Only 1 fully public implementation exists.",
                CONTENT_LEFT, Inches(5.8), CONTENT_WIDTH, Inches(0.5),
                size=CAPTION_SIZE, italic=True, color=MID_GRAY)


def slide_03_solution(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "The Solution")
    add_footer(slide)

    add_textbox(slide,
                "A domain-specific AI agent that understands part numbers, model compatibility, "
                "symptoms, and installation — not a generic chatbot wrapper",
                CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(0.8),
                size=BODY_SIZE, bold=False, color=DARK)

    capabilities = [
        ("\u2699  Part Lookup", " — hybrid semantic + keyword search"),
        ("\u2705  Compatibility Check", " — 162,976 model mappings"),
        ("\U0001F50D  Symptom Diagnosis", " — 159 repair guides + 51 blogs"),
        ("\U0001F6E0  Installation Guidance", " — difficulty, video, steps"),
        ("\U0001F6E1  Scope Guardrails", " — 3-layer defense"),
        ("\U0001F3A4  Voice Input", " — Web Speech API"),
    ]
    add_bullet_list(slide, capabilities, CONTENT_LEFT, Inches(2.4), CONTENT_WIDTH, Inches(3.5),
                    size=Pt(17))

    add_textbox(slide, "Live demo: website-chat-agent.vercel.app",
                CONTENT_LEFT, Inches(5.8), CONTENT_WIDTH, Inches(0.5),
                size=SMALL_SIZE, bold=True, color=TEAL)


def slide_04_welcome(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Interface: Welcome & Starter Prompts")
    add_footer(slide)

    img_w, img_h = Inches(9.0), Inches(5.0)
    img_left = (SLIDE_W - img_w) // 2
    add_image_or_placeholder(slide, "starter-prompts.png", img_left, Inches(1.3), img_w, img_h)

    add_textbox(slide, "Guided entry points for common tasks — reduces blank-page anxiety",
                CONTENT_LEFT, Inches(6.5), CONTENT_WIDTH, Inches(0.5),
                size=CAPTION_SIZE, italic=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_05_conversation(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Interface: Multi-Turn Conversation")
    add_footer(slide)

    img_w, img_h = Inches(9.0), Inches(5.0)
    img_left = (SLIDE_W - img_w) // 2
    add_image_or_placeholder(slide, "hero-conversation.png", img_left, Inches(1.3), img_w, img_h)

    add_textbox(slide,
                "Rich components stream in real-time: DiagnosisCards, ProductCards, quick-reply suggestions",
                CONTENT_LEFT, Inches(6.5), CONTENT_WIDTH, Inches(0.5),
                size=CAPTION_SIZE, italic=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_06_product_card(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Interface: Rich Product Cards")
    add_footer(slide)

    # Image on left
    add_image_or_placeholder(slide, "product-card.png",
                             Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5))

    # Bullets on right
    items = [
        "Product image, name, brand, PS number",
        "Price with stock status badge",
        "5-star rating with review count",
        "Installation difficulty indicator",
        '"Fixes:" symptom tags',
        "Direct link to PartSelect.com product page",
    ]
    add_bullet_list(slide, items, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5),
                    size=Pt(16))


def slide_07_diagnosis_card(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Interface: Symptom Diagnosis")
    add_footer(slide)

    # Image on left
    add_image_or_placeholder(slide, "diagnosis-card.png",
                             Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5))

    items = [
        "Ranked causes with likelihood (High / Medium / Low)",
        "Recommended replacement parts per cause",
        "Part thumbnails with prices",
        "Follow-up question buttons for refinement",
        "Sources linked to repair guides",
    ]
    add_bullet_list(slide, items, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5),
                    size=Pt(16))


def slide_08_compat_voice_guard(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Interface: Compatibility, Voice, & Guardrails")
    add_footer(slide)

    # Three images in a row
    img_w = Inches(3.6)
    img_h = Inches(3.8)
    gap = Inches(0.5)
    start_left = Inches(0.7)

    files = ["compatibility-badge.png", "voice-input.png", "guardrail.png"]
    captions = [
        "Verified compatible badge\nwith part/model numbers",
        "Mic button for hands-free\nqueries mid-repair",
        "Polite redirect for\nout-of-scope appliances",
    ]

    for i, (fname, cap) in enumerate(zip(files, captions)):
        left = int(start_left + i * (img_w + gap))
        add_image_or_placeholder(slide, fname, left, Inches(1.3), img_w, img_h)
        # Caption
        txBox = slide.shapes.add_textbox(left, Inches(5.3), img_w, Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = cap
        _set_font(run, size=CAPTION_SIZE, color=DARK, italic=True)


def slide_09_architecture(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Architecture: Custom Agent Loop")
    add_footer(slide)

    diagram = (
        "User  \u2192  Next.js Frontend  \u2192  SSE Stream  \u2192  FastAPI Backend\n"
        "                                                    \u2193\n"
        "                                          Agent Loop (while True)\n"
        "                                   LLM decides \u2192 Tool call \u2192 Result \u2192 LLM decides \u2192 ...\n"
        "                                                    \u2193\n"
        "                                          Tool Layer (5 tools)\n"
        "                                                    \u2193\n"
        "                                       Data Layer (ChromaDB + JSON)"
    )
    add_code_block(slide, diagram, Inches(0.7), Inches(1.5), Inches(11.9), Inches(3.8))

    add_textbox(slide,
                "~50 lines of Python. No LangChain. No LlamaIndex. Fully debuggable.",
                CONTENT_LEFT, Inches(5.7), CONTENT_WIDTH, Inches(0.6),
                size=BODY_SIZE, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)


def slide_10_why_custom(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Why Not LangChain?")
    add_footer(slide)

    items = [
        ("Debuggable — ", "The entire orchestration is ~50 lines. Every tool call, every LLM decision "
                          "is logged and visible. No framework internals to trace through."),
        ("Multi-intent — ", '"Is this part compatible AND how do I install it?" chains '
                            "check_compatibility \u2192 get_installation_guide naturally. "
                            "A state machine breaks here."),
        ("Signal — ", "For a case study, custom code signals stronger engineering than wrapping a "
                      "framework. Trivial to explain in a Q&A."),
    ]
    add_bullet_list(slide, items, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(4.5),
                    size=Pt(17), spacing=Pt(18))


def slide_11_five_tools(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Five Domain-Specific Tools")
    add_footer(slide)

    headers = ["Tool", "What It Does"]
    rows = [
        ["search_parts", "Hybrid semantic + keyword search, 3-tier fallback (scoped \u2192 all \u2192 not found)"],
        ["check_compatibility", "Exact lookup against 162,976 model mappings. Honest \"not in data\" vs false negatives"],
        ["diagnose_symptom", "Symptom \u2192 ranked causes \u2192 recommended parts. Searches repair guides + blogs"],
        ["get_installation_guide", "Steps, difficulty (1\u20135), time estimate, YouTube video link"],
        ["get_product_details", "Full product data for rich card rendering"],
    ]
    add_table(slide, headers, rows, Inches(0.7), Inches(1.4), Inches(11.9), Inches(0.55))

    add_callout_box(slide,
                    'Every tool has a required "reasoning" parameter — chain-of-thought that improves '
                    "accuracy and creates an audit trail",
                    Inches(0.7), Inches(5.2), Inches(11.9), Inches(0.7))


def slide_12_walkthrough(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Walkthrough: Multi-Step Query")
    add_footer(slide)

    add_textbox(slide, 'User: "My dishwasher won\'t drain. I have model WDT780SAEM1."',
                CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(0.5),
                size=Pt(17), bold=True, color=DARK, italic=True)

    steps = [
        '1.  Agent calls diagnose_symptom(symptom="won\'t drain", appliance_type="dishwasher")',
        "2.  Returns 6 ranked causes (Piston & Nut Assembly, Drain Pump, Check Valve...)",
        '3.  Agent calls search_parts(query="dishwasher drain pump")',
        "4.  Returns matching parts with prices and ratings",
        '5.  Agent calls check_compatibility(part_number="PS...", model_number="WDT780SAEM1")',
        '6.  Returns "Verified Compatible \u2713"',
        "7.  Final response: diagnosis + compatible part recommendation + install link",
    ]
    add_bullet_list(slide, steps, CONTENT_LEFT, Inches(2.1), CONTENT_WIDTH, Inches(3.8),
                    size=Pt(15), spacing=Pt(6))

    add_textbox(slide,
                "All 3 tool calls happen in a single turn — the LLM chains them automatically",
                CONTENT_LEFT, Inches(5.8), CONTENT_WIDTH, Inches(0.5),
                size=CAPTION_SIZE, italic=True, color=MID_GRAY)


def slide_13_llm_choice(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "LLM Choice: Gemini 2.5 Flash")
    add_footer(slide)

    headers = ["Model", "TTFT", "Throughput", "Cost (input/1M)"]
    rows = [
        ["Gemini 2.5 Flash", "~250ms", "~250 tok/s", "$0.15"],
        ["GPT-4.1 Mini", "~800ms", "~62 tok/s", "$0.40"],
        ["Claude Haiku 4.5", "~600ms", "~79 tok/s", "$0.80"],
    ]
    add_table(slide, headers, rows, Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.55))

    add_textbox(slide,
                "For a customer support chatbot, perceived speed matters enormously. "
                "250ms TTFT vs 800ms is the difference between 'instant' and 'loading...'",
                CONTENT_LEFT, Inches(3.8), CONTENT_WIDTH, Inches(0.8),
                size=BODY_SIZE, color=DARK)

    add_textbox(slide,
                "Fallback: If Gemini is down \u2192 GPT-4.1 Mini (same function-calling format)",
                CONTENT_LEFT, Inches(5.0), CONTENT_WIDTH, Inches(0.5),
                size=SMALL_SIZE, italic=True, color=MID_GRAY)


def slide_14_data_pipeline(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Data: 4,170 Parts, Scraped with Playwright")
    add_footer(slide)

    add_textbox(slide, "Playwright with real Chrome bypasses Akamai anti-bot protection",
                CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(0.5),
                size=BODY_SIZE, bold=True, color=TEAL)

    headers = ["Dataset", "Count", "Source"]
    rows = [
        ["Parts (fully enhanced)", "4,170", "Playwright (Chrome)"],
        ["Model Mappings", "162,976", "AJAX scroll expansion"],
        ["Repair Guides", "159", "22 generic + 101 brand + 36 how-to"],
        ["Blog Articles", "51", "86,587 words"],
        ["Vectors in ChromaDB", "~5,107", "Parts + knowledge"],
    ]
    add_table(slide, headers, rows, Inches(0.7), Inches(2.2), Inches(11.9), Inches(0.50))

    add_callout_box(slide,
                    "Initial Firecrawl scrape: 101 parts, 30 models/part.  \u2192  "
                    "Playwright scrape: 4,170 parts, 429 models/part.",
                    Inches(0.7), Inches(5.3), Inches(11.9), Inches(0.7))


def slide_15_embedding(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Embedding Strategy: One Part = One Vector")
    add_footer(slide)

    # Left column — Traditional RAG (crossed out)
    add_textbox(slide, "\u2718  Traditional RAG",
                Inches(0.7), CONTENT_TOP, Inches(5.5), Inches(0.5),
                size=Pt(20), bold=True, color=RED)
    trad_items = [
        "Chunk raw HTML into 512-token blocks",
        "Navigation noise, image URLs dilute quality",
        "Part data split across multiple chunks",
    ]
    add_bullet_list(slide, trad_items, Inches(0.7), Inches(2.1), Inches(5.5), Inches(2.0),
                    size=Pt(15))

    # Right column — Our Approach
    add_textbox(slide, "\u2714  Our Approach",
                Inches(7.0), CONTENT_TOP, Inches(5.5), Inches(0.5),
                size=Pt(20), bold=True, color=TEAL)
    our_items = [
        "Construct a search-optimized document per part",
        "One part = one vector, no chunking artifacts",
        "Metadata filtering (appliance_type, brand, chunk_type)",
    ]
    add_bullet_list(slide, our_items, Inches(7.0), Inches(2.1), Inches(5.5), Inches(2.0),
                    size=Pt(15))

    # Code example
    code = (
        'Refrigerator Door Shelf Bin by Whirlpool. PS11752778 / WPW10321304.\n'
        '$36.18, 4.9 stars (351 reviews). In stock.\n'
        'Fixes symptoms: Door won\'t close, Ice buildup, Leaking.'
    )
    add_code_block(slide, code, Inches(0.7), Inches(4.5), Inches(11.9), Inches(1.5))


def slide_16_guardrails(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Three-Layer Guardrails")
    add_footer(slide)

    layers = [
        ("1. System Prompt", "(primary)",
         "Instructs LLM to specialize in fridge/dishwasher, politely redirect off-topic"),
        ("2. Rule-Based Classifier", "(zero latency, zero cost)",
         "Keyword/regex runs BEFORE the LLM. Off-topic \u2192 canned redirect, never hits the LLM"),
        ("3. Metadata Filtering", "(data-level)",
         'ChromaDB tier-1 queries restricted to appliance_type IN ["refrigerator", "dishwasher"]. '
         "Even if LLM ignores prompt, data layer enforces scope"),
    ]

    top = Inches(1.5)
    for title, badge, desc in layers:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.7), top, Inches(11.9), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF0, 0xF7, 0xF7)
        box.line.color.rgb = TEAL
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        run_t = p.add_run()
        run_t.text = title + "  "
        _set_font(run_t, size=Pt(17), bold=True, color=TEAL)
        run_b = p.add_run()
        run_b.text = badge
        _set_font(run_b, size=SMALL_SIZE, color=MID_GRAY, italic=True)

        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        run_d = p2.add_run()
        run_d.text = desc
        _set_font(run_d, size=Pt(15), color=DARK)

        top += Inches(1.75)


def slide_17_extensibility(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Extensibility: Built to Grow")
    add_footer(slide)

    scenarios = [
        ("New Appliance (washing machines) \u2192 1\u20132 hours\n",
         'Add "washer" to tier-1 filter, scrape repair guides, update system prompt. '
         "No re-scraping parts, no re-embedding, no code changes to agent loop."),
        ("New Tool (order tracking) \u2192 Half a day\n",
         "Define Python function, register with Gemini, add SSE event type, build frontend component."),
        ("Swap LLM (Gemini \u2192 GPT-4.1) \u2192 One config value\n",
         "Agent loop is model-agnostic, tools use standard JSON schema."),
    ]

    top = Inches(1.5)
    for header, body in scenarios:
        txBox = slide.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_WIDTH, Inches(1.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run_h = p.add_run()
        run_h.text = header
        _set_font(run_h, size=Pt(17), bold=True, color=TEAL)
        run_b = p.add_run()
        run_b.text = body
        _set_font(run_b, size=Pt(15), color=DARK)
        top += Inches(1.6)

    add_textbox(slide,
                "Why easy: We scraped ALL appliance types upfront. Scope filtering happens at query time, not ingestion.",
                CONTENT_LEFT, Inches(6.0), CONTENT_WIDTH, Inches(0.5),
                size=CAPTION_SIZE, italic=True, color=MID_GRAY)


def slide_18_tech_stack(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Tech Stack")
    add_footer(slide)

    headers = ["Layer", "Technology", "Why"]
    rows = [
        ["Frontend", "Next.js 16 + React 19 + Tailwind 4", "Server components, streaming, brand styling"],
        ["Backend", "FastAPI + Python 3.13", "Async-native, SSE streaming, fast startup"],
        ["LLM", "Gemini 2.5 Flash", "Fastest TTFT, highest throughput, lowest cost"],
        ["Embeddings", "Gemini Embedding (768d)", "Task-type optimization, multimodal-ready"],
        ["Vector DB", "ChromaDB", "Zero-config, metadata filtering"],
        ["Scraping", "Playwright (Chrome)", "Bypasses Akamai, scrolls AJAX pagination"],
        ["Streaming", "SSE", "Works through Vercel edge/CDN, auto-reconnect"],
        ["Deploy", "Vercel + Docker", "Free tier, edge CDN"],
    ]
    add_table(slide, headers, rows, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.50))


def slide_19_roadmap(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "What's Next")
    add_footer(slide)

    items = [
        "Image upload for visual part identification (Gemini Embedding supports multimodal vectors)",
        "TTS read-aloud for hands-free installation guidance",
        "Order tracking tool integration",
        "Part comparison side-by-side",
        "Production: Redis sessions, Pinecone vector DB, Kubernetes scaling, Vertex AI for SLAs",
    ]
    add_bullet_list(slide, items, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(4.5),
                    size=BODY_SIZE, spacing=Pt(14))


def slide_20_closing(prs):
    slide = new_slide(prs)

    # Teal bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    # Title
    add_textbox(slide, "PartSelect AI Chat Agent",
                Inches(1.5), Inches(1.2), Inches(10), Inches(0.8),
                size=Pt(36), bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)

    # 4 criteria checkmarks
    criteria = [
        "\u2713  Interface Design — Rich cards, voice input, guided flows, PartSelect branded",
        "\u2713  Agentic Architecture — Custom 50-line agent loop, 5 tools, reasoning audit trail",
        "\u2713  Extensibility — New appliance in 2 hours, new tool in half a day, swap LLM in 1 config",
        "\u2713  Query Accuracy — 162,976 models, 159 guides, hybrid search, honest confidence",
    ]
    add_bullet_list(slide, criteria,
                    Inches(1.2), Inches(2.4), Inches(10.8), Inches(3.0),
                    size=Pt(16), spacing=Pt(12))

    # Links
    add_textbox(slide, "Live Demo: website-chat-agent.vercel.app",
                Inches(1.5), Inches(5.0), Inches(10), Inches(0.4),
                size=SMALL_SIZE, color=TEAL, alignment=PP_ALIGN.CENTER, bold=True)

    add_textbox(slide, "Built by Aryaman Gupta for Instalily AI",
                Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
                size=BODY_SIZE, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_welcome(prs)
    slide_05_conversation(prs)
    slide_06_product_card(prs)
    slide_07_diagnosis_card(prs)
    slide_08_compat_voice_guard(prs)
    slide_09_architecture(prs)
    slide_10_why_custom(prs)
    slide_11_five_tools(prs)
    slide_12_walkthrough(prs)
    slide_13_llm_choice(prs)
    slide_14_data_pipeline(prs)
    slide_15_embedding(prs)
    slide_16_guardrails(prs)
    slide_17_extensibility(prs)
    slide_18_tech_stack(prs)
    slide_19_roadmap(prs)
    slide_20_closing(prs)

    prs.save(str(OUTPUT_PATH))
    print(f"Saved {len(prs.slides)} slides to {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
